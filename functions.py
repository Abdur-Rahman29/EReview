# functions.py
import difflib
from docx import Document as DocxDocument
import pdfplumber 
import streamlit as st
import re
from pptx import Presentation
def load_documents_from_files(files):
    documents = {}
    for file in files:
        if file.name.endswith(".txt"):
            documents[file.name] = file.getvalue().decode("utf-8")
        elif file.name.endswith(".docx"):
            doc = DocxDocument(file)
            full_text = [para.text for para in doc.paragraphs]
            documents[file.name] = '\n'.join(full_text)
        elif file.name.endswith(".pdf"):
            with pdfplumber.open(file) as pdf:
                full_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                documents[file.name] = '\n'.join(full_text)
        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            documents[file.name] = '\n'.join(text)
    return documents

def extract_changed_code(old_code, new_code):
    """Extract modified lines and clean them for review."""
    old_code_lines = old_code.splitlines(keepends=True)
    new_code_lines = new_code.splitlines(keepends=True)

    # Calculate the diff
    diff = list(difflib.unified_diff(old_code_lines, new_code_lines, lineterm=''))
    return diff  # Return the raw diff output

def parse_diff(old_code, new_code, diff):
    changes = []
    old_line_no = 0
    new_line_no = 0
    old_code_lines = old_code.splitlines()
    new_code_lines = new_code.splitlines()

    # Track removed line waiting for corresponding addition
    pending_removal = None

    for line in diff:
        if line.startswith('-'):
            # Line removed from old_code
            old_line_no += 1
            change_content = line[1:].strip()  # Content of the removed line
            if change_content and change_content not in ('--', '++'):
                pending_removal = {
                    "Old Line": old_line_no,
                    "Change": change_content
                }
        elif line.startswith('+'):
            # Line added to new_code
            new_line_no += 1
            change_content = line[1:].strip()  # Content of the added line
            if change_content and change_content not in ('--', '++'):
                if pending_removal:
                    # Handle modification
                    changes.append({
                        "Old Line": pending_removal["Old Line"],
                        "New Line": new_line_no,
                        "Change": f'{pending_removal["Change"]} --> {change_content}',
                        "Type": "Modified"
                    })
                    pending_removal = None  # Clear pending removal after matching
                else:
                    # No pending removal, this is a new addition
                    changes.append({
                        "Old Line": '',
                        "New Line": new_line_no,
                        "Change": change_content,
                        "Type": "Added"
                    })
        elif line.startswith(' '):
            # Unchanged line, increment both counters
            old_line_no += 1
            new_line_no += 1
            # If there's a pending removal that wasn't matched, mark it as removed
            if pending_removal:
                changes.append({
                    "Old Line": pending_removal["Old Line"],
                    "New Line": '',
                    "Change": pending_removal["Change"],
                    "Type": "Removed"
                })
                pending_removal = None  # Clear the pending removal

    # If there's a pending removal left at the end, append it as a removal
    if pending_removal:
        changes.append({
            "Old Line": pending_removal["Old Line"],
            "New Line": '',
            "Change": pending_removal["Change"],
            "Type": "Removed"
        })

    # Adjust line numbers based on actual content
    for change in changes:
        if change["Type"] == "Removed" or change["Type"] == "Modified":
            # Find the actual line number in old_code for removal or modification
            for idx, old_line in enumerate(old_code_lines):
                if old_line.strip() == change["Change"].split(" --> ")[0].strip():
                    change["Old Line"] = idx + 1  # Convert to 1-based index
                    break
        if change["Type"] == "Added" or change["Type"] == "Modified":
            # Find the actual line number in new_code for addition or modification
            for idx, new_line in enumerate(new_code_lines):
                if new_line.strip() == change["Change"].split(" --> ")[-1].strip():
                    change["New Line"] = idx + 1  # Convert to 1-based index
                    break
    return changes
    
def compare_code(old_code, new_code):
    old_code_lines = old_code.splitlines(keepends=True)
    new_code_lines = new_code.splitlines(keepends=True)
    
    diff = list(difflib.unified_diff(old_code_lines, new_code_lines, lineterm=''))
    
    return ("Modified", diff) if len(diff) > 2 else ("New", new_code_lines)

def display_error_tabs(code_file_content, client,org_standards_documents):
    error_counts={}

    with st.expander("Errors In New Code", expanded=False):
        # Initialize error counts
        syntax_error_count = 0
        runtime_error_count = 0
        logical_error_count = 0

        # Collect Syntax Errors
        if code_file_content:
            syntax_prompt = (
                "Identify any syntax errors in the following code. "
                "If none are found, respond with 'No errors'. "
                "For each error, use this format: "
                "Line No: [line number], Error: [error line] in 1st line\n"
                "Description: [error description] in 2nd line\n"
                "Suggestion: [correction] in 3rd line.\n"
                "Do not omit comments in your evaluation. Don't give anything else.\n"
                f"Code:\n{code_file_content}"
            )
            syntax_response = client.chat.completions.create(
                messages=[{"role": "user", "content": syntax_prompt}],
                model="llama3-8b-8192",
            )
            syntax_errors = syntax_response.choices[0].message.content
            syntax_error_count = len(re.findall(r"Line No: \d+", syntax_errors))

        # Collect Runtime Errors
        
            runtime_prompt = (
                "Identify any runtime errors in the following code. "
                "If none are found, respond with 'No errors'. "
                "For each error, use this format: "
                "Line No: [line number], Error: [error line] in 1st line\n"
                "Description: [error description] in 2nd line\n"
                "Suggestion: [correction] in 3rd line.\n"
                "Do not omit comments in your evaluation. Don't give anything else.\n"
                f"Code:\n{code_file_content}"
            )
            runtime_response = client.chat.completions.create(
                messages=[{"role": "user", "content": runtime_prompt}],
                model="llama3-8b-8192",
            )
            runtime_errors = runtime_response.choices[0].message.content
            runtime_error_count = len(re.findall(r"Line No: \d+", runtime_errors))

        # Collect Logical Errors
       
            logical_prompt = (
                "Identify any logical errors in the following code. "
                "If none are found, respond with 'No errors'. "
                "For each error, use this format: "
                "Line No: [line number], Error: [error line] in 1st line\n"
                "Description: [error description] in 2nd line\n"
                "Suggestion: [correction] in 3rd line.\n"
                "Do not omit comments in your evaluation. Don't give anything else.\n"
                f"Code:\n{code_file_content}"
            )

            logical_response = client.chat.completions.create(
                messages=[{"role": "user", "content": logical_prompt}],
                model="llama3-8b-8192",
            )
            logical_errors = logical_response.choices[0].message.content
            logical_error_count = len(re.findall(r"Line No: \d+", logical_errors))

            validation_prompt = (
                "Only if Language is HTML or CSS: Correctly Identify any Validation Errors in the following code, else: Display No errors"
                "For each error, use this format: "
                "Line No: [line number], Error: [error line] in 1st line\n"
                "Description: [error description] in 2nd line\n"
                "Suggestion: [correction] in 3rd line.\n"
                "Do not omit comments in your evaluation. Don't give anything else.\n"
                f"Code:\n{code_file_content}"
            )
            validation_response = client.chat.completions.create(
                messages=[{"role": "user", "content": validation_prompt}],
                model="llama3-8b-8192",
            )
            validation_errors = validation_response.choices[0].message.content
            validation_error_count = len(re.findall(r"Line No: \d+", validation_errors))

            compilation_prompt = (
                "Only if Language is Java: Correctly Identify any Compile-time in the following code, else: Display No errors"
                "For each error, use this format: "
                "Line No: [line number], Error: [error line] in 1st line\n"
                "Description: [error description] in 2nd line\n"
                "Suggestion: [correction] in 3rd line.\n"
                "Do not omit comments in your evaluation. Don't give anything else.\n"
                f"Code:\n{code_file_content}"
            )
            compilation_response = client.chat.completions.create(
                messages=[{"role": "user", "content": compilation_prompt}],
                model="llama3-8b-8192",
            )
            compilation_errors = compilation_response.choices[0].message.content
            compilation_error_count = len(re.findall(r"Line No: \d+", compilation_errors))

        # Create tabs with error counts
        tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
            f"Syntax Errors ({syntax_error_count})",
            f"Runtime Errors ({runtime_error_count})",
            f"Logical Errors ({logical_error_count})",
            f"Validation Errors ({validation_error_count})",
            f"Compilation Errors ({compilation_error_count})",
            "Suggested Code"
        ])
        error_counts={
          'Syntax Errors': syntax_error_count,
          'Runtime Errors': runtime_error_count,
          'Logical Errors': logical_error_count,
          'Validation Errors': validation_error_count,
          'Compilation Errors': compilation_error_count
        }
        with tab1:
            st.write("### Syntax Errors")
            st.write(syntax_errors)

        with tab2:
            st.write("### Runtime Errors")
            st.write(runtime_errors)

        with tab3:
            st.write("### Logical Errors")
            st.write(logical_errors)
        with tab4:
            st.write("### Validation Errors")
            st.write(validation_errors)
        with tab5:
            st.write("### Compilation Errors")
            st.write(compilation_errors)
        with tab6:
            st.write("### Suggestions")
            if code_file_content:
                suggested_prompt = (
                    "Improve the errors found in Code and provide suggestions"
                    f"Code:\n{code_file_content}"
                    
                )
                suggested_response = client.chat.completions.create(
                    messages=[{"role": "user", "content": suggested_prompt}],
                    model="llama3-8b-8192",
                )
                st.write(suggested_response.choices[0].message.content)
        return error_counts

def calculate_score(org_std_text, new_code,client):
    # Define prompts based on organizational code standards
    explain_prompt = (
                        "Give only Score out of 10 of following code based on Organisation Standards:\n"
                        f"Code: {new_code}"
                        f"Organisation Standards: {org_std_text}"
                    )
            
    explain_response = client.chat.completions.create(
                        messages=[{"role": "user", "content": explain_prompt}],
                        model="llama3-8b-8192",
                    )
                    
    explanation = explain_response.choices[0].message.content
    
    # Assuming the response is a score as a float, parse it
    score_match = re.search(r'\d+(\.\d+)?', explanation)  # This will match a number (integer or float)

    if score_match:
        score = float(score_match.group(0))  # Extract the score and convert it to a float
    else:
        st.error("Failed to parse the score from the response.")
        return None, None

    # Display color based on score
    if score >= 9:
        color = "green"
        message = "Excellent"
    elif score >= 7:
        color = "yellow"
        message = "Good"
    elif score >= 4:
        color = "orange"
        message = "Average"
    else:
        color = "red"
        message = "Poor"

    # Display the score with color
    score_colored = f'<span style="color:{color}; font-size: 20px;">{score}</span>'
    message_colored = f'<span style="color:{color}; font-size: 20px;">{message}</span>'
    st.subheader('Overall score based on standards')
    st.markdown(f'Overall Score: {score_colored} - {message_colored}', unsafe_allow_html=True)

    return explanation  # Ensure both are returned

error_weights = {
    "Syntax Errors": 3,
    "Run-time Errors": 4,
    "Logical Errors": 2,
    "Validation Errors": 1,
    "Compile-time Errors": 4
}

# Function to calculate total severity score from the dictionary {error_type: count}
def calculate_severity(error_count):
    total_score = 0
    
    # Iterate over the error_count dictionary
    for error_type, count in error_count.items():
        # Check if the error_type exists in the predefined weights
        if error_type in error_weights:
            # Multiply the count by the weight of the error type and add it to the total score
            total_score += count * error_weights[error_type]
    
    return total_score

# Determine severity level based on the total score
def determine_severity_from_score(total_score):
    if total_score == 0:
        color = "darkgreen"
        message = "No errors, perfect!"
    elif total_score <= 10:
        color = "lightgreen"
        message =  "Low severity"
    elif 11 <= total_score <= 20:
        color = "yellow"
        message = "Medium severity"
    elif 21 <= total_score <= 30:
        color = "orange"
        message = "High severity"
    else:
        color = "red"
        message = "Critical severity"
    

    return color, message
def detect_vulnerabilities(code):
    vulnerabilities = []

    # Example checks for vulnerabilities
    if "eval(" in code or "exec(" in code:
        vulnerabilities.append("Avoid using eval() or exec() as they can lead to code injection vulnerabilities.")

    if re.search(r"(\bSELECT\b|\bUPDATE\b|\bDELETE\b|\bINSERT\b)\s+\w+\s+FROM\s+\w+\s*;", code, re.IGNORECASE):
        vulnerabilities.append("Possible SQL Injection vulnerability detected. Use parameterized queries instead.")

    # Check for common XSS patterns
    if re.search(r"<script.*?>", code, re.IGNORECASE):
        vulnerabilities.append("Potential Cross-Site Scripting (XSS) vulnerability detected. Avoid using inline scripts.")

    # Check for hardcoded credentials (this is a simplistic check)
    if re.search(r"(password|secret|api_key)\s*=\s*['\"].*['\"]", code, re.IGNORECASE):
        vulnerabilities.append("Hardcoded credentials detected. Consider using environment variables or a secure vault.")

    return vulnerabilities

def severity(error_count,total_score,code_to_review):
    total_score = calculate_severity(error_count)  # Calculate total severity score
    color, severity = determine_severity_from_score(total_score)  # Determine the severity level

    # Display the score with color
    st.markdown(f"### Severity Analysis")
    st.write(f"Total Errors: {sum(error_count.values())}")
    total_score_colored = f'<span style="color:{color}; font-size: 20px;">{total_score}</span>'
    message_colored = f'<span style="color:{color}; font-size: 20px;">{severity}</span>'
    st.markdown(f'Total Severity Score: {total_score_colored } - Severity: { message_colored }', unsafe_allow_html=True)

    vulnerabilities_found = detect_vulnerabilities(code_to_review)
    # Display vulnerabilities and improvements
    if vulnerabilities_found:
        st.subheader("Detected Vulnerabilities")
        for vulnerability in vulnerabilities_found:
            st.markdown(f"- {vulnerability}")
    else:
        st.success("No vulnerabilities detected.")  
